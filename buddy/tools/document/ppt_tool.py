from __future__ import annotations

from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

from buddy.prompts.ppt_prompts import PPT_TOOL_PROMPT

_TOOL = "ppt"

# ── constants ─────────────────────────────────────────────────────────────────

# Slide dimensions in EMU (914400 per inch)
_ASPECT_DIMS: Dict[str, Tuple[int, int]] = {
    "16:9":     (12192000, 6858000),
    "4:3":      (9144000,  6858000),
    "portrait": (6858000,  9144000),
}

_THEMES: Dict[str, Dict[str, str]] = {
    "light":     {"bg": "FFFFFF", "text": "1A1A1A", "accent": "2563EB"},
    "dark":      {"bg": "1A1A2E", "text": "E2E8F0", "accent": "6366F1"},
    "minimal":   {"bg": "F8FAFC", "text": "0F172A", "accent": "0EA5E9"},
    "corporate": {"bg": "0F3460", "text": "FFFFFF", "accent": "E94560"},
}

_FONT_PT: Dict[str, int] = {
    "small": 14, "medium": 20, "large": 28, "heading": 36,
}

# Zone layout: (left_frac, top_frac, width_frac, height_frac)
_ZONES: Dict[str, Tuple[float, float, float, float]] = {
    "title":          (0.05, 0.03, 0.90, 0.18),
    "subtitle":       (0.10, 0.55, 0.80, 0.20),
    "body":           (0.05, 0.22, 0.90, 0.72),
    "left":           (0.03, 0.22, 0.46, 0.72),
    "right":          (0.51, 0.22, 0.46, 0.72),
    "top_left":       (0.03, 0.22, 0.46, 0.35),
    "top_right":      (0.51, 0.22, 0.46, 0.35),
    "bottom_left":    (0.03, 0.57, 0.46, 0.37),
    "bottom_right":   (0.51, 0.57, 0.46, 0.37),
    "full":           (0.00, 0.00, 1.00, 1.00),
    "title_center":   (0.10, 0.28, 0.80, 0.28),
    "subtitle_center":(0.15, 0.58, 0.70, 0.18),
    "section_center": (0.10, 0.32, 0.80, 0.36),
}

# Image sub-positions within their parent zone: (left_frac, top_frac, width_frac, height_frac)
_IMG_POSITIONS: Dict[str, Optional[Tuple[float, float, float, float]]] = {
    "full":         None,
    "center":       (0.25, 0.10, 0.50, 0.80),
    "left":         (0.00, 0.10, 0.50, 0.80),
    "right":        (0.50, 0.10, 0.50, 0.80),
    "top_left":     (0.00, 0.00, 0.50, 0.50),
    "top_right":    (0.50, 0.00, 0.50, 0.50),
    "bottom_left":  (0.00, 0.50, 0.50, 0.50),
    "bottom_right": (0.50, 0.50, 0.50, 0.50),
}

_IMG_SIZES: Dict[str, float] = {
    "small": 0.25, "medium": 0.50, "large": 0.75, "fill": 1.00,
}

_VALID_OPS = "add_slide, update_slide, delete_slide, reorder, set_background, add_element"


# ── helpers ───────────────────────────────────────────────────────────────────

def _resolve(raw: str) -> str:
    import os
    p = os.path.expanduser(os.path.expandvars(raw.strip()))
    if not os.path.isabs(p):
        p = os.path.join(os.path.expanduser("~"), p)
    return p


def _ok(**kw: Any) -> Dict[str, Any]:
    return {"STATUS": "success", **kw}


def _err(msg: str, **kw: Any) -> Dict[str, Any]:
    return {"STATUS": "failed", "TOOL": _TOOL, "ERROR": msg, **kw}


def _needs_confirm(preview: str, **kw: Any) -> Dict[str, Any]:
    return {
        "STATUS": "failed", "TOOL": _TOOL,
        "NEEDS_CONFIRMATION": True,
        "PREVIEW": preview,
        "NOTE": "Call again with confirmed=true after user approves.",
        **kw,
    }


def _hex_rgb(h: str) -> Any:
    from pptx.dml.color import RGBColor
    h = h.lstrip("#")[:6].upper()
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _compute_zone(zone: str, slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    lf, tf, wf, hf = _ZONES.get(zone, _ZONES["body"])
    return (int(slide_w * lf), int(slide_h * tf), int(slide_w * wf), int(slide_h * hf))


def _detect_aspect_ratio(prs: Any) -> str:
    w, h = prs.slide_width, prs.slide_height
    if w is None or h is None:
        return "16:9"
    ratio = w / h
    if abs(ratio - 12192000 / 6858000) < 0.05:
        return "16:9"
    if abs(ratio - 9144000 / 6858000) < 0.05:
        return "4:3"
    if w < h:
        return "portrait"
    return "16:9"


def _get_blank_layout(prs: Any) -> Any:
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[-1]


# ── slide rendering ───────────────────────────────────────────────────────────

def _apply_bg(slide: Any, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_rgb(color_hex)


def _apply_para_style(para: Any, style: Dict) -> None:
    from pptx.enum.text import PP_ALIGN
    _align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    para.alignment = _align.get((style or {}).get("align", "left"), PP_ALIGN.LEFT)


def _apply_run_style(run: Any, style: Dict, default_color: str) -> None:
    from pptx.util import Pt
    s = style or {}
    run.font.size = Pt(_FONT_PT.get(s.get("size", "medium"), 20))
    run.font.bold = bool(s.get("bold", False))
    run.font.italic = bool(s.get("italic", False))
    run.font.color.rgb = _hex_rgb(s.get("color") or default_color)


def _add_text_box(
    slide: Any, text: str, l: int, t: int, w: int, h: int,
    style: Dict, default_color: str,
) -> None:
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    _apply_para_style(para, style)
    run = para.add_run()
    run.text = text
    _apply_run_style(run, style, default_color)


def _add_bullets_box(
    slide: Any, items: List[str], l: int, t: int, w: int, h: int,
    style: Dict, default_color: str,
) -> None:
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _apply_para_style(para, style)
        run = para.add_run()
        run.text = f"• {item}"
        _apply_run_style(run, style, default_color)


def _add_solid_rect(slide: Any, l: int, t: int, w: int, h: int, hex_color: str) -> None:
    """Add a colored rectangle, supporting 8-char hex (#RRGGBBAA) for alpha."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    h_str = hex_color.lstrip("#")
    rgb_hex = h_str[:6].upper()
    alpha_val = 100000
    if len(h_str) == 8:
        alpha_byte = int(h_str[6:8], 16)
        alpha_val = int(alpha_byte / 255 * 100000)

    shape = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_rgb(rgb_hex)

    if alpha_val < 100000:
        sp = shape._element
        solid_fill = sp.find(".//" + qn("a:solidFill"))
        if solid_fill is not None:
            srgb_clr = solid_fill.find(qn("a:srgbClr"))
            if srgb_clr is not None:
                for old in srgb_clr.findall(qn("a:alpha")):
                    srgb_clr.remove(old)
                alpha_elem = etree.SubElement(srgb_clr, qn("a:alpha"))
                alpha_elem.set("val", str(alpha_val))


def _add_image_element(
    slide: Any, elem: Dict, zone_emu: Tuple[int, int, int, int],
    slide_w: int, slide_h: int,
) -> Dict[str, Any]:
    from pptx.util import Emu

    raw_path = str(elem.get("path") or "").strip()
    if not raw_path:
        return _err("image element missing 'path'")
    img_path = Path(_resolve(raw_path))
    if not img_path.exists():
        return _err(f"IMAGE_NOT_FOUND: {img_path}")

    position = str(elem.get("position") or "center").lower()
    size_key = str(elem.get("size") or "large").lower()
    background = elem.get("background")
    overlay = elem.get("overlay")

    zl, zt, zw, zh = zone_emu

    if position == "full":
        img_l, img_t, img_w, img_h = 0, 0, slide_w, slide_h
    else:
        pos_fracs = _IMG_POSITIONS.get(position, _IMG_POSITIONS["center"])
        pfl, pft, pfw, pfh = pos_fracs
        area_l = zl + int(zw * pfl)
        area_t = zt + int(zh * pft)
        area_w = int(zw * pfw)
        area_h = int(zh * pfh)

        size_frac = _IMG_SIZES.get(size_key, 0.75)
        img_w = int(area_w * size_frac)

        try:
            from PIL import Image as PILImage
            with PILImage.open(str(img_path)) as pil_img:
                pw, ph = pil_img.size
                aspect = pw / ph if ph > 0 else 4.0 / 3.0
        except Exception:
            aspect = 4.0 / 3.0

        img_h = int(img_w / aspect)
        img_l = area_l + (area_w - img_w) // 2
        img_t = area_t + (area_h - img_h) // 2

    if background:
        _add_solid_rect(slide, img_l, img_t, img_w, img_h, background)

    try:
        slide.shapes.add_picture(str(img_path), Emu(img_l), Emu(img_t), Emu(img_w), Emu(img_h))
    except Exception as e:
        return _err(f"Failed to add image '{img_path.name}': {e}")

    if overlay:
        _add_solid_rect(slide, img_l, img_t, img_w, img_h, overlay)

    return _ok()


def _add_table_element(
    slide: Any, elem: Dict, l: int, t: int, w: int, h: int, theme: Dict,
) -> None:
    from pptx.util import Emu, Pt

    headers = list(elem.get("headers") or [])
    rows = list(elem.get("rows") or [])
    if not headers and not rows:
        return

    n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    n_rows = len(rows) + (1 if headers else 0)

    tbl = slide.shapes.add_table(n_rows, n_cols, Emu(l), Emu(t), Emu(w), Emu(h)).table

    row_offset = 0
    if headers:
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(hdr)
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = _hex_rgb("FFFFFF")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_rgb(theme.get("accent", "2563EB"))
        row_offset = 1

    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            val = row[ci] if ci < len(row) else ""
            cell = tbl.cell(ri + row_offset, ci)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(val) if val is not None else ""
            run.font.size = Pt(13)
            run.font.color.rgb = _hex_rgb(theme.get("text", "1A1A1A"))


def _render_element(
    slide: Any, elem: Dict, zone_emu: Tuple[int, int, int, int],
    theme: Dict, default_color: str, slide_w: int, slide_h: int,
) -> Dict[str, Any]:
    elem_type = str(elem.get("type") or "text").lower()
    style = elem.get("style") or {}
    l, t, w, h = zone_emu

    if elem_type == "text":
        value = str(elem.get("value") or elem.get("content") or "")
        _add_text_box(slide, value, l, t, w, h, style, default_color)
    elif elem_type == "bullets":
        items = [str(x) for x in (elem.get("items") or [])]
        _add_bullets_box(slide, items, l, t, w, h, style, default_color)
    elif elem_type == "image":
        result = _add_image_element(slide, elem, zone_emu, slide_w, slide_h)
        if result.get("STATUS") == "failed":
            return result
    elif elem_type == "table":
        _add_table_element(slide, elem, l, t, w, h, theme)
    else:
        return _err(f"Unknown element type: '{elem_type}'. Valid: text, bullets, image, table")

    return _ok(type=elem_type)


def _render_elements_in_zone(
    slide: Any, elements: List[Dict], zone_emu: Tuple[int, int, int, int],
    theme: Dict, default_color: str, slide_w: int, slide_h: int,
) -> Optional[Dict[str, Any]]:
    if not elements:
        return None
    zl, zt, zw, zh = zone_emu
    n = len(elements)
    slot_h = zh // n
    for i, elem in enumerate(elements):
        slot = (zl, zt + i * slot_h, zw, slot_h)
        result = _render_element(slide, elem, slot, theme, default_color, slide_w, slide_h)
        if result.get("STATUS") == "failed":
            return result
    return None


def _populate_slide(
    slide: Any, spec: Dict, theme: Dict, slide_w: int, slide_h: int,
) -> Optional[Dict[str, Any]]:
    layout = str(spec.get("layout") or "title_content").lower()
    bg_color = spec.get("background") or theme.get("bg", "FFFFFF")
    default_color = theme.get("text", "1A1A1A")

    _apply_bg(slide, bg_color)

    title = spec.get("title")
    subtitle = spec.get("subtitle")
    content = list(spec.get("content") or [])
    content_left = list(spec.get("content_left") or [])
    content_right = list(spec.get("content_right") or [])
    notes_text = spec.get("notes")

    if title:
        if layout == "title_slide":
            zone_name = "title_center"
            title_style: Dict = {"size": "heading", "bold": True, "align": "center"}
        elif layout == "section_header":
            zone_name = "section_center"
            title_style = {"size": "heading", "bold": True, "align": "center"}
        else:
            zone_name = "title"
            title_style = {"size": "large", "bold": True, "align": "left"}
        tl, tt, tw, th = _compute_zone(zone_name, slide_w, slide_h)
        _add_text_box(slide, title, tl, tt, tw, th, title_style, default_color)

    if subtitle and layout == "title_slide":
        sl, st, sw, sh = _compute_zone("subtitle_center", slide_w, slide_h)
        _add_text_box(slide, subtitle, sl, st, sw, sh,
                      {"size": "medium", "align": "center"}, default_color)

    if layout == "two_column":
        err = _render_elements_in_zone(
            slide, content_left, _compute_zone("left", slide_w, slide_h),
            theme, default_color, slide_w, slide_h,
        )
        if err:
            return err
        err = _render_elements_in_zone(
            slide, content_right, _compute_zone("right", slide_w, slide_h),
            theme, default_color, slide_w, slide_h,
        )
        if err:
            return err
    elif layout == "blank":
        err = _render_elements_in_zone(
            slide, content, _compute_zone("full", slide_w, slide_h),
            theme, default_color, slide_w, slide_h,
        )
        if err:
            return err
    else:
        err = _render_elements_in_zone(
            slide, content, _compute_zone("body", slide_w, slide_h),
            theme, default_color, slide_w, slide_h,
        )
        if err:
            return err

    if notes_text:
        try:
            slide.notes_slide.notes_text_frame.text = str(notes_text)
        except Exception:
            pass

    return None


def _clear_slide_content(slide: Any) -> None:
    from pptx.oxml.ns import qn
    sp_tree = slide.shapes._spTree
    removable = {qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"), qn("p:grpSp"), qn("p:cxnSp")}
    for child in list(sp_tree):
        if child.tag in removable:
            sp_tree.remove(child)


def _extract_slide_info(slide: Any, idx: int) -> Dict[str, Any]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import os
    import hashlib

    spec: Dict[str, Any] = {"slide_number": idx + 1}
    content: List[Dict[str, Any]] = []
    texts: List[List[str]] = []

    try:
        layout_name = slide.slide_layout.name
        spec["layout"] = layout_name
    except Exception:
        spec["layout"] = "unknown"

    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            spec["notes"] = notes
    except Exception:
        pass

    for shape in slide.shapes:
        try:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                ext = shape.image.ext
                h = hashlib.md5(image_bytes).hexdigest()
                img_dir = os.path.expanduser("~/.buddy/ppt_media")
                os.makedirs(img_dir, exist_ok=True)
                img_path = os.path.join(img_dir, f"{h}.{ext}")
                
                if not os.path.exists(img_path):
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                
                content.append({
                    "type": "image",
                    "path": img_path,
                    "position": "center",
                    "size": "large"
                })
            
            elif getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text_frame.text.strip())
                    rows.append(row_data)
                
                if rows:
                    if len(rows) > 1:
                        content.append({"type": "table", "headers": rows[0], "rows": rows[1:]})
                    else:
                        content.append({"type": "table", "headers": [], "rows": rows})
                        
            elif getattr(shape, "has_text_frame", False):
                para_texts = []
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        para_texts.append(t)
                
                if para_texts:
                    texts.append(para_texts)
                    
        except Exception:
            # Safely ignore shapes we cannot extract
            pass

    if texts:
        spec["title"] = " ".join(texts[0])
        for para_texts in texts[1:]:
            if len(para_texts) > 1:
                content.append({"type": "bullets", "items": para_texts})
            else:
                content.append({"type": "text", "value": " ".join(para_texts)})
                
    if content:
        spec["content"] = content
    else:
        spec["content"] = []

    return spec


# ── edit op handlers ──────────────────────────────────────────────────────────

def _op_add_slide(
    prs: Any, op: Dict, theme: Dict, slide_w: int, slide_h: int,
) -> Dict[str, Any]:
    spec = op.get("slide")
    if not isinstance(spec, dict):
        return {"op": "add_slide", "STATUS": "failed", "ERROR": "'slide' object is required"}

    blank = _get_blank_layout(prs)
    new_slide = prs.slides.add_slide(blank)
    err = _populate_slide(new_slide, spec, theme, slide_w, slide_h)
    if err:
        return {"op": "add_slide", **err}

    at = op.get("at")
    total = len(prs.slides)
    if at is not None:
        target_idx = max(0, min(int(at) - 1, total - 1))
        current_idx = total - 1
        if target_idx != current_idx:
            xml_slides = prs.slides._sldIdLst
            elem = xml_slides[current_idx]
            xml_slides.remove(elem)
            xml_slides.insert(target_idx, elem)

    return {"op": "add_slide", "STATUS": "success", "TOTAL_SLIDES": len(prs.slides)}


def _op_update_slide(
    prs: Any, op: Dict, theme: Dict, slide_w: int, slide_h: int,
) -> Dict[str, Any]:
    slide_number = op.get("slide_number")
    spec = op.get("slide")
    if slide_number is None:
        return {"op": "update_slide", "STATUS": "failed", "ERROR": "slide_number is required"}
    if not isinstance(spec, dict):
        return {"op": "update_slide", "STATUS": "failed", "ERROR": "'slide' object is required"}

    idx = int(slide_number) - 1
    total = len(prs.slides)
    if idx < 0 or idx >= total:
        return {"op": "update_slide", "STATUS": "failed",
                "ERROR": f"slide_number {slide_number} out of range — presentation has {total} slides"}

    slide = prs.slides[idx]
    
    # Deep extract existing slide into a perfect JSON spec
    existing_spec = _extract_slide_info(slide, idx)
    merged_spec = dict(existing_spec)
    
    # Perform a smart union: overwrite old elements with explicitly provided new elements
    for k, v in spec.items():
        if v is not None:
            merged_spec[k] = v

    _clear_slide_content(slide)
    err = _populate_slide(slide, merged_spec, theme, slide_w, slide_h)
    if err:
        return {"op": "update_slide", **err}

    return {"op": "update_slide", "STATUS": "success", "SLIDE_NUMBER": int(slide_number)}


def _op_delete_slide(prs: Any, op: Dict) -> Dict[str, Any]:
    slide_number = op.get("slide_number")
    if slide_number is None:
        return {"op": "delete_slide", "STATUS": "failed", "ERROR": "slide_number is required"}
    if not op.get("confirmed"):
        return {
            "op": "delete_slide", "STATUS": "failed",
            "NEEDS_CONFIRMATION": True,
            "PREVIEW": f"Will permanently delete slide {slide_number}.",
            "NOTE": 'Add "confirmed": true inside this op to proceed.',
        }

    idx = int(slide_number) - 1
    total = len(prs.slides)
    if idx < 0 or idx >= total:
        return {"op": "delete_slide", "STATUS": "failed",
                "ERROR": f"slide_number {slide_number} out of range — presentation has {total} slides"}

    xml_slides = prs.slides._sldIdLst
    sldId_elem = xml_slides[idx]
    r_id = sldId_elem.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    xml_slides.remove(sldId_elem)
    if r_id:
        try:
            prs.part.drop_rel(r_id)
        except Exception:
            pass

    return {"op": "delete_slide", "STATUS": "success", "DELETED_SLIDE_NUMBER": int(slide_number)}


def _op_reorder(prs: Any, op: Dict) -> Dict[str, Any]:
    from_num = op.get("from")
    to_num = op.get("to")
    if from_num is None or to_num is None:
        return {"op": "reorder", "STATUS": "failed", "ERROR": "'from' and 'to' are both required (1-based)"}

    total = len(prs.slides)
    from_idx = int(from_num) - 1
    to_idx = int(to_num) - 1

    for label, idx in [("from", from_idx), ("to", to_idx)]:
        if idx < 0 or idx >= total:
            return {"op": "reorder", "STATUS": "failed",
                    "ERROR": f"{label}={idx+1} out of range — presentation has {total} slides"}

    xml_slides = prs.slides._sldIdLst
    elem = xml_slides[from_idx]
    xml_slides.remove(elem)
    xml_slides.insert(to_idx, elem)

    return {"op": "reorder", "STATUS": "success", "FROM": int(from_num), "TO": int(to_num)}


def _op_set_background(prs: Any, op: Dict) -> Dict[str, Any]:
    slide_number = op.get("slide_number")
    color = str(op.get("color") or "").strip()
    if not color:
        return {"op": "set_background", "STATUS": "failed", "ERROR": "'color' is required (hex string)"}
    if slide_number is None:
        return {"op": "set_background", "STATUS": "failed", "ERROR": "'slide_number' is required"}

    if str(slide_number).lower() == "all":
        for s in prs.slides:
            _apply_bg(s, color)
        return {"op": "set_background", "STATUS": "success", "AFFECTED": "all", "COLOR": color}

    idx = int(slide_number) - 1
    total = len(prs.slides)
    if idx < 0 or idx >= total:
        return {"op": "set_background", "STATUS": "failed",
                "ERROR": f"slide_number {slide_number} out of range — presentation has {total} slides"}

    _apply_bg(prs.slides[idx], color)
    return {"op": "set_background", "STATUS": "success",
            "SLIDE_NUMBER": int(slide_number), "COLOR": color}


def _op_add_element(
    prs: Any, op: Dict, theme: Dict, slide_w: int, slide_h: int,
) -> Dict[str, Any]:
    slide_number = op.get("slide_number")
    elem = op.get("element")
    if slide_number is None:
        return {"op": "add_element", "STATUS": "failed", "ERROR": "slide_number is required"}
    if not isinstance(elem, dict):
        return {"op": "add_element", "STATUS": "failed", "ERROR": "'element' object is required"}

    idx = int(slide_number) - 1
    total = len(prs.slides)
    if idx < 0 or idx >= total:
        return {"op": "add_element", "STATUS": "failed",
                "ERROR": f"slide_number {slide_number} out of range — presentation has {total} slides"}

    slide = prs.slides[idx]
    zone_emu = _compute_zone("body", slide_w, slide_h)
    result = _render_element(
        slide, elem, zone_emu, theme, theme.get("text", "1A1A1A"), slide_w, slide_h,
    )
    if result.get("STATUS") == "failed":
        return {"op": "add_element", **result}

    return {"op": "add_element", "STATUS": "success", "SLIDE_NUMBER": int(slide_number)}


# ── tool ──────────────────────────────────────────────────────────────────────

class PptTool:
    tool_name = _TOOL
    version = "1.0.0"

    def get_info(self) -> Dict[str, Any]:
        return {
            "name": self.tool_name,
            "version": self.version,
            "description": (
                "WHEN: any operation on a .pptx presentation — creating slides, editing content.\n\n"
                "FUNCTIONS:\n"
                "  create(path, slides[], aspect_ratio?, theme?)     — new presentation; slides: [{layout, title, content[], ...}]\n"
                "  read(path)                                         — returns SLIDE_COUNT, ASPECT_RATIO, SLIDES[]\n"
                "  edit(path, operations[], theme?)                   — batch ops: add_slide, update_slide, delete_slide, reorder, set_background, add_element\n\n"
                "CHAIN: always call read before edit to get correct slide_numbers (1-based, shift after delete/reorder).\n"
                "NOT: .ppt (legacy) | .doc/.docx → word | .xlsx → excel | .pdf → pdf | export/convert → converter"
            ),
            "prompt": PPT_TOOL_PROMPT,
        }

    async def execute(
        self,
        function: str,
        arguments: Dict[str, Any],
        on_progress: Optional[Callable] = None,
        goal: str = "",
        brain: Any = None,
        **_: Any,
    ) -> Dict[str, Any]:
        fn = str(function or "").strip().lower()
        if fn == "create":
            return self._create(arguments)
        if fn == "read":
            return self._read(arguments)
        if fn == "edit":
            return self._edit(arguments)
        return _err(f"Unknown function: {function!r}. Must be: create, read, edit")

    # ── create ────────────────────────────────────────────────────────────────

    def _create(self, args: Dict[str, Any]) -> Dict[str, Any]:
        raw = str(args.get("path") or "").strip()
        if not raw:
            return _err("ppt.create: 'path' is required — provide an absolute .pptx path")
        if not raw.lower().endswith(".pptx"):
            return _err("ppt.create: path must end in .pptx")

        slides_data = args.get("slides")
        if not slides_data or not isinstance(slides_data, list):
            return _err("ppt.create: 'slides' is required and must be a non-empty list")

        path = _resolve(raw)
        p = Path(path)

        if p.exists() and not args.get("confirmed"):
            return _needs_confirm(
                f"File already exists: {path}\n"
                "Setting confirmed=true will overwrite it and all its slides."
            )

        try:
            from pptx import Presentation
        except ImportError:
            return _err("python-pptx not installed. Run: pip install python-pptx")

        aspect_ratio = str(args.get("aspect_ratio") or "16:9").strip()
        if aspect_ratio not in _ASPECT_DIMS:
            aspect_ratio = "16:9"
        slide_w, slide_h = _ASPECT_DIMS[aspect_ratio]

        theme_name = str(args.get("theme") or "light").lower()
        theme = _THEMES.get(theme_name, _THEMES["light"])

        try:
            prs = Presentation()
            prs.slide_width = slide_w
            prs.slide_height = slide_h
            blank = _get_blank_layout(prs)

            for i, spec in enumerate(slides_data):
                if not isinstance(spec, dict):
                    return _err(f"slides[{i}] must be an object with layout, title, content, etc.")
                new_slide = prs.slides.add_slide(blank)
                err = _populate_slide(new_slide, spec, theme, slide_w, slide_h)
                if err:
                    return _err(f"slides[{i}]: {err.get('ERROR', 'render failed')}")

            p.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(p))
            return _ok(PATH=str(p), SLIDES=len(slides_data),
                       ASPECT_RATIO=aspect_ratio, THEME=theme_name)

        except PermissionError:
            return _err(f"Permission denied: {path}")
        except Exception as e:
            return _err(f"Failed to create presentation: {e}")

    # ── read ──────────────────────────────────────────────────────────────────

    def _read(self, args: Dict[str, Any]) -> Dict[str, Any]:
        raw = str(args.get("path") or "").strip()
        if not raw:
            return _err("ppt.read: 'path' is required — provide an absolute .pptx path")

        path = _resolve(raw)
        p = Path(path)
        if not p.exists():
            return _err(f"ppt.read: file not found: {path} — use fs_browse.find to locate it")
        if not p.is_file():
            return _err(f"ppt.read: path is a directory, not a file: {path}")

        try:
            from pptx import Presentation
        except ImportError:
            return _err("python-pptx not installed. Run: pip install python-pptx")

        try:
            prs = Presentation(str(p))
            aspect_ratio = _detect_aspect_ratio(prs)
            slides_info = [_extract_slide_info(s, i) for i, s in enumerate(prs.slides)]
            return _ok(
                PATH=str(p),
                SLIDE_COUNT=len(prs.slides),
                ASPECT_RATIO=aspect_ratio,
                SLIDES=slides_info,
            )
        except PermissionError:
            return _err(f"Permission denied: {path}")
        except Exception as e:
            return _err(f"Failed to read presentation: {e}")

    # ── edit ──────────────────────────────────────────────────────────────────

    def _edit(self, args: Dict[str, Any]) -> Dict[str, Any]:
        raw = str(args.get("path") or "").strip()
        if not raw:
            return _err("ppt.edit: 'path' is required — provide an absolute .pptx path")

        operations = args.get("operations")
        if not operations or not isinstance(operations, list):
            return _err("ppt.edit: 'operations' is required and must be a non-empty list")

        path = _resolve(raw)
        p = Path(path)
        if not p.exists():
            return _err(f"ppt.edit: file not found: {path} — use fs_browse.find to locate it")

        try:
            from pptx import Presentation
        except ImportError:
            return _err("python-pptx not installed. Run: pip install python-pptx")

        try:
            prs = Presentation(str(p))
        except PermissionError:
            return _err(f"Permission denied: {path}")
        except Exception as e:
            return _err(f"Failed to open presentation: {e}")

        slide_w = prs.slide_width or _ASPECT_DIMS["16:9"][0]
        slide_h = prs.slide_height or _ASPECT_DIMS["16:9"][1]
        theme_name = str(args.get("theme") or "light").lower()
        theme = _THEMES.get(theme_name, _THEMES["light"])

        results: List[Dict[str, Any]] = []
        for op in operations:
            if not isinstance(op, dict):
                results.append(_err(f"Each operation must be an object, got: {type(op).__name__}"))
                break
            result = self._apply_op(prs, op, theme, slide_w, slide_h)
            results.append(result)
            if result.get("STATUS") == "failed":
                break

        all_ok = all(r.get("STATUS") == "success" for r in results)

        if all_ok:
            try:
                prs.save(str(p))
            except PermissionError:
                return _err(
                    f"Permission denied when saving: {path} — file may be open in another app.",
                    RESULTS=results, SAVED=False,
                )
            except Exception as e:
                return _err(f"Operations succeeded but save failed: {e}", RESULTS=results, SAVED=False)

        out: Dict[str, Any] = {
            "STATUS": "success" if all_ok else "failed",
            "PATH": str(p),
            "OPERATIONS_TOTAL": len(operations),
            "OPERATIONS_APPLIED": len(results),
            "RESULTS": results,
            "SAVED": all_ok,
        }
        if not all_ok:
            failed = next((r for r in results if r.get("STATUS") == "failed"), {})
            out["ERROR"] = failed.get("ERROR") or "An operation failed — file was not modified."
        return out

    def _apply_op(
        self, prs: Any, op: Dict, theme: Dict, slide_w: int, slide_h: int,
    ) -> Dict[str, Any]:
        op_type = str(op.get("op") or "").strip().lower()

        try:
            if op_type == "add_slide":
                return _op_add_slide(prs, op, theme, slide_w, slide_h)
            if op_type == "update_slide":
                return _op_update_slide(prs, op, theme, slide_w, slide_h)
            if op_type == "delete_slide":
                return _op_delete_slide(prs, op)
            if op_type == "reorder":
                return _op_reorder(prs, op)
            if op_type == "set_background":
                return _op_set_background(prs, op)
            if op_type == "add_element":
                return _op_add_element(prs, op, theme, slide_w, slide_h)

            return {
                "op": op_type or "(empty)", "STATUS": "failed",
                "ERROR": f"Unknown op: '{op_type}'. Valid ops: {_VALID_OPS}",
                "VALID_OPS": _VALID_OPS,
            }
        except Exception as e:
            return {"op": op_type, "STATUS": "failed", "ERROR": str(e)}



# ── registry ──────────────────────────────────────────────────────────────────

TOOL_NAME = "ppt"
TOOL_CLASS = PptTool


def get_tool() -> PptTool:
    return PptTool()
